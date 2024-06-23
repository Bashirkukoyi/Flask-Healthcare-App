from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add a title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Survey Data Analysis"
subtitle.text = "Generated using Python"

# Add a slide for Age vs Income chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Top 10 Ages with the Highest Total Income"
img_path = 'ages_with_highest_income.png'
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8.5))

# Add a slide for Gender Distribution chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Gender Distribution Across Spending Categories"
img_path = 'gender_distribution_spending.png'
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8.5))

# Save the presentation
prs.save('survey_data_analysis.pptx')
